"""
Description: [One-line technical summary of the file's primary responsibility]
Project: PhD Thesis - Physics-Informed Protection for IBR-Dominated Microgrids
Author: Dr. Anoop Eluvathingal (Senior Scientist)
Logic Layer: Presentation Delivery
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_slide(prs, title, points):
    slide_layout = prs.slide_layouts[1] # Bullet point layout
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.word_wrap = True
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.space_after = Pt(14)

def create_presentation():
    prs = Presentation()

    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.shapes.placeholders[1]
    title.text = "Physics-Informed Hybrid Protection for IBR-Dominated Microgrids"
    subtitle.text = "Dr. Anoop Eluvathingal\nSenior Scientist | IIT Madras PhD Candidate"

    # Slide 2: The Magnitude Problem
    add_slide(prs, "The Problem: IBR Current-Limiting", [
        "Inverter-Based Resources (IBRs) suppress fault current magnitudes.",
        "Traditional ANSI 50/51 relays fail to detect High Impedance Faults (HIF).",
        "Result: 'Invisible' faults leading to catastrophic equipment failure."
    ])

    # Slide 3: Method A Failure (The "Naive" Approach)
    add_slide(prs, "Method A Failure: Magnitude Bias", [
        "Initial Model: MLP trained on raw current magnitude (ib).",
        "Performance: 100% Sensitivity in lab, but 99.9% Nuisance Tripping in audit.",
        "Finding: -29.9% Security Margin during normal load transients.",
        "Conclusion: Magnitude is a false correlate for IBR protection."
    ])

    # Slide 4: Method B Innovation (The Entropy Fix)
    add_slide(prs, "Method B: Residual Entropy Analysis", [
        "Innovation: Isolating Arcing Chaos from the Fundamental 50Hz Sine.",
        "Physics: i_residual = i_total - i_rolling_mean.",
        "Primary Features: Noise Density (Sigma) and Ripple Peak-to-Peak.",
        "Benefit: Magnitude-invariant fault discrimination."
    ])

    # Slide 5: The 8-Scenario Performance Matrix
    add_slide(prs, "Validation: The 8-Scenario Matrix", [
        "Sensitivity: Detected 0.6 p.u. faults (HIF Case A).",
        "Security: Blocked 1.5 p.u. Motor Starts (Case B).",
        "Robustness: Immune to Cap-Switching and Inverter Harmonics.",
        "System Security Margin: Successfully recovered to >40%."
    ])

    # Slide 6: Coordination & Selectivity
    add_slide(prs, "Selectivity & IEC 61850 Integration", [
        "Case E: Successful discrimination of adjacent-feeder through-faults.",
        "Communication: GOOSE-assisted coordination logic.",
        "Result: Maximum microgrid availability during external disturbances."
    ])

    # Slide 7: Hardware & Industrial Roadmap
    add_slide(prs, "Implementation: Edge-AI Requirements", [
        "Hardware: ARM Cortex-M7 / TI C2000 DSP for 100μs inference.",
        "Process: Real-time digital filtering + MLP execution.",
        "Integration: ANSI 21LB (Load Blinder) augmented IED."
    ])

    # Slide 8: Conclusion
    add_slide(prs, "Conclusion: Doctorate Readiness", [
        "Redefined protection boundaries for IBR-dominated grids.",
        "Proven transition from Stochastic AI to Physics-Informed Logic.",
        "Field-ready roadmap for next-gen Microgrid Protection."
    ])

    prs.save('Anoop_PhD_Synopsis.pptx')
    print("SUCCESS: Anoop_PhD_Synopsis.pptx generated.")

if __name__ == "__main__":
    create_presentation()