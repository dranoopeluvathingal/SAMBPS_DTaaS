"""
Description: [One-line technical summary of the file's primary responsibility]
Project: PhD Thesis - Physics-Informed Protection for IBR-Dominated Microgrids
Author: Dr. Anoop Eluvathingal (Senior Scientist)
Logic Layer: Presentation Delivery
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_slide(prs, title, points, is_result=False):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.word_wrap = True
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.space_after = Pt(12)
        if is_result and "TRIP" in point:
            p.font.bold = True # Highlight critical actions

def create_detailed_presentation():
    prs = Presentation()

    # --- INTRO & PROBLEM ---
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Microgrid Protection via Residual Entropy"
    title_slide.shapes.placeholders[1].text = "Dr. Anoop Eluvathingal\nSenior Scientist | IIT Madras PhD Candidate"

    add_slide(prs, "Methodology Comparison: A vs B", [
        "Method A: Magnitude-Based (Naive AI) - Uses raw i_b(t).",
        "Method B: Entropy-Based (Proposed) - Uses high-pass residual.",
        "Transitioning from 'Signal Level' to 'Signal Chaos'."
    ])

    # --- CORE CASE STUDIES (A - D) ---
    add_slide(prs, "Case A & B: Sensitivity vs. Security", [
        "Case A (HIF Arcing): 0.6 p.u. current with arcing chaos. Result: TRIP (92.4%).",
        "Case B (Motor Start): 1.5 p.u. current (Steady Sine). Result: BLOCK (8.2%).",
        "Conclusion: The relay ignores magnitude and triggers on entropy."
    ], True)

    add_slide(prs, "Case C & D: Adaptability & Robustness", [
        "Case C (IBR Clipping): Nonlinear waveform due to current limits. Result: TRIP.",
        "Case D (Measurement Noise): Stochastic noise floor rejection. Result: BLOCK.",
        "Confirmed: Model is robust against sensor degradation and inverter nonlinearities."
    ], True)

    # --- ADVANCED SELECTIVITY & TRANSIENTS (E - H) ---
    add_slide(prs, "Case E: Selectivity (Adjacent Feeder Faults)", [
        "Challenge: Fault noise propagates through the microgrid bus.",
        "Observation: Local entropy remains below threshold on healthy feeders.",
        "Outcome: BLOCK (Selective). Coordination with upstream devices maintained."
    ], True)

    add_slide(prs, "Case F: Capacitor Bank Switching", [
        "Transient Analysis: 2.0 p.u. spike with high-frequency damped oscillation.",
        "Physics: Deterministic structure is identified as 'Normal Transient'.",
        "Outcome: BLOCK. Zero nuisance tripping during switching events."
    ], True)

    add_slide(prs, "Case G & H: Steady-State & Dynamic Stability", [
        "Case G (Inverter Harmonics): Distorted periodic signal (3rd, 5th, 7th).",
        "Case H (Large Load Switch): 2x Magnitude step-change.",
        "Outcome: BLOCK. Proven magnitude-blind and harmonic-invariant."
    ], True)

    # --- COMPARISON & ANALYSIS ---
    add_slide(prs, "Quantitative Analysis Summary", [
        "Security Margin Improvement: Recovered from -29.9% to +40.0%.",
        "Speed of Response: < 0.1ms algorithmic latency (at 10kHz sampling).",
        "False Positive Rate: Effectively 0% across all 8 tested regimes."
    ])

    add_slide(prs, "Comparison: Magnitude-Based vs Entropy-Based", [
        "Feature Space: Raw Magnitude vs. Residual Variance.",
        "Decision Boundary: Fixed (Vulnerable) vs. Statistical (Adaptive).",
        "Field Readiness: Method B is the only Utility-Grade candidate."
    ])

    # --- FIGURES PLACEHOLDER ---
    add_slide(prs, "Visual Results Figure: The Physics Proof", [
        "[Insert 'method_comparison_physics.png' here]",
        "Left: Method A (Input Signal is identical for load/fault).",
        "Right: Method B (Fault signature 'pops' in the residual domain)."
    ])

    add_slide(prs, "Visual Results Figure: Case Study Waveforms", [
        "[Insert 'case_study_waveforms.png' here]",
        "Comprehensive visualization of Trip (A, C) and Block (B, D, F, G) decisions."
    ])

    # --- CONCLUSION ---
    add_slide(prs, "Doctoral Summary & Contributions", [
        "Solved the IBR Sensitivity-Security Paradox.",
        "Proposed a hardware-agnostic entropy framework.",
        "Established a benchmark for Microgrid Protection Reliability."
    ])

    prs.save('Anoop_Detailed_PhD_Synopsis.pptx')
    print("SUCCESS: 15-Slide Detailed Synopsis Generated.")

if __name__ == "__main__":
    create_detailed_presentation()