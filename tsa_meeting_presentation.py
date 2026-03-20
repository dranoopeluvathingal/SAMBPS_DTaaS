from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # Initialize presentation
    prs = Presentation()
    
    # --- Set Slide Dimensions (Standard 4:3 or 16:9) ---
    # Beamer default is often 4:3, but 16:9 is modern. We'll use 16:9.
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # --- Define Colors (IIT Madras/LaTeX Theme) ---
    maroon = RGBColor(128, 0, 0)
    black = RGBColor(0, 0, 0)
    dark_gray = RGBColor(64, 64, 64)
    white = RGBColor(255, 255, 255)
    
    # --- Helper Function: Add Header and Footer to Content Slides ---
    def add_slide_frame(slide, title_text, slide_num):
        # Top Header Bar
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
        header.fill.solid()
        header.fill.fore_color.rgb = maroon
        header.line.color.rgb = maroon
        
        # Title Text
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), prs.slide_width - Inches(1), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(32)
        p.font.color.rgb = white
        p.font.bold = True
        
        # Footer Bar
        footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.5), prs.slide_width, Inches(0.5))
        footer.fill.solid()
        footer.fill.fore_color.rgb = maroon
        footer.line.color.rgb = maroon
        
        # Footer Content (Left: Name/Topic, Right: Slide Num)
        txBox_f1 = slide.shapes.add_textbox(Inches(0.2), prs.slide_height - Inches(0.45), Inches(6), Inches(0.4))
        tf_f1 = txBox_f1.text_frame
        p_f1 = tf_f1.paragraphs[0]
        p_f1.text = "Anoop V Eluvathingal | TSA Meeting | IIT Madras"
        p_f1.font.size = Pt(12)
        p_f1.font.color.rgb = white
        
        txBox_f2 = slide.shapes.add_textbox(prs.slide_width - Inches(1.5), prs.slide_height - Inches(0.45), Inches(1), Inches(0.4))
        tf_f2 = txBox_f2.text_frame
        p_f2 = tf_f2.paragraphs[0]
        p_f2.text = f"{slide_num}/12"
        p_f2.font.size = Pt(12)
        p_f2.font.color.rgb = white
        p_f2.alignment = PP_ALIGN.RIGHT
        
        return slide
    
    # --- Helper Function: Add Bulleted List ---
    def add_bullet_list(slide, items, top=Inches(1.5), left=Inches(1.0), width=Inches(11.333), height=Inches(5.0)):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(items):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(24)
            p.font.color.rgb = black
            p.level = 0
            p.space_after = Pt(14)
            
    # ==========================================
    # SLIDE 1: Title Slide (LaTeX Beamer Style)
    # ==========================================
    slide_layout = prs.slide_layouts[6] # Blank
    slide1 = prs.slides.add_slide(slide_layout)
    
    # Optional: Add IITM Logo placeholder if you want
    # slide1.shapes.add_picture('iitm_logo.png', Inches(5.66), Inches(0.5), width=Inches(2.0))
    
    # Title
    txBox = slide1.shapes.add_textbox(Inches(1), Inches(2.0), prs.slide_width - Inches(2), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Advanced Protection for Distributed\nPower Systems and Microgrids"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = maroon
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle / Meeting Details
    txBox2 = slide1.shapes.add_textbox(Inches(2), Inches(4.0), prs.slide_width - Inches(4), Inches(2))
    tf2 = txBox2.text_frame
    
    lines = [
        ("Thesis Submission Approval Meeting", 24, True),
        ("by", 18, False),
        ("Anoop V Eluvathingal", 28, True),
        ("Doctor of Philosophy", 20, False),
        ("Dept. of Electrical Engineering", 20, False),
        ("IIT Madras", 20, False),
        ("June 2026", 18, False)
    ]
    
    for i, (text, size, is_bold) in enumerate(lines):
        p = tf2.add_paragraph() if i > 0 else tf2.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = is_bold
        p.alignment = PP_ALIGN.CENTER
        if text == "by" or text == "IIT Madras":
            p.space_after = Pt(10)
    
    # Footer Bar for Title Slide
    footer = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.5), prs.slide_width, Inches(0.5))
    footer.fill.solid()
    footer.fill.fore_color.rgb = maroon
    footer.line.color.rgb = maroon


    # ==========================================
    # SLIDE 2: Outline
    # ==========================================
    slide2 = add_slide_frame(prs.slides.add_slide(slide_layout), "Outline", 2)
    bullets2 = [
        "Introduction & Background",
        "Research Motivation",
        "Literature Review Summary",
        "Research Objectives",
        "Key Contributions",
        "Thesis Organization",
        "List of Publications",
        "Conclusions"
    ]
    add_bullet_list(slide2, bullets2)

    # ==========================================
    # SLIDE 3: Background & Context
    # ==========================================
    slide3 = add_slide_frame(prs.slides.add_slide(slide_layout), "Background and Context", 3)
    bullets3 = [
        "Evolving topological frameworks of Active Distribution Networks.",
        "Complex fault signature characterization of Inverter-Based Resources (IBRs).",
        "Stringent regulatory landscape and interconnection compliance (IEEE 1547 / IEC 61727).",
        "Critical impact of reduced rotational inertia on microgrid fault recovery and stability."
    ]
    add_bullet_list(slide3, bullets3)

    # ==========================================
    # SLIDE 4: Research Motivation
    # ==========================================
    slide4 = add_slide_frame(prs.slides.add_slide(slide_layout), "Research Motivation", 4)
    bullets4 = [
        "Need for physics-based modelling of DERs to understand transient responses under unbalanced faults.",
        "Desensitisation and mal-operation risks inherent in legacy protection relaying principles.",
        "Coordination challenges in bi-directional, low-inertia networks.",
        "Limitations of standard Interconnection Protection Relay logics for IBRs.",
        "The necessity of AI/ML approaches for fault detection in non-deterministic microgrid environments."
    ]
    add_bullet_list(slide4, bullets4)

    # ==========================================
    # SLIDE 5: Literature Review Summary
    # ==========================================
    slide5 = add_slide_frame(prs.slides.add_slide(slide_layout), "State-of-the-Art & Limitations", 5)
    bullets5 = [
        "Industrial maturity of dynamic settings and adaptive principles in modern IEDs.",
        "Identified limitations of Sequence-Component Based Directional Discrimination.",
        "Adoption of IEC 61850 based GOOSE messaging in industry.",
        "Emerging trends: Digital Twins and Physics-Informed Neural Networks (PINNs) applied to power systems."
    ]
    add_bullet_list(slide5, bullets5)

    # ==========================================
    # SLIDE 6: Research Objectives
    # ==========================================
    slide6 = add_slide_frame(prs.slides.add_slide(slide_layout), "Research Objectives", 6)
    bullets6 = [
        "Objective 1: Develop Fault-Resilient Adaptive Interconnection Protection Logics for IBR-Dominated Distribution Networks.",
        "Objective 2: Synthesize and implement Hierarchical Microgrid Protection architectures.",
        "Objective 3: Explore and design Data-Driven Paradigms (AI/ML) for advanced Microgrid protection."
    ]
    add_bullet_list(slide6, bullets6)

    # ==========================================
    # SLIDE 7: Key Contributions (1/2)
    # ==========================================
    slide7 = add_slide_frame(prs.slides.add_slide(slide_layout), "Key Contributions", 7)
    bullets7 = [
        "C1: Development of adaptive protection logics for Multifunctional Interconnection Protection Relays.",
        "Formulation of instantaneous time-domain sampling logics for frequency-independent fault detection.",
        "Validation against deterministic legacy relays across extreme scenarios (highly resistive faults, bidirectional infeeds, islanding)."
    ]
    add_bullet_list(slide7, bullets7)

    # ==========================================
    # SLIDE 8: Key Contributions (2/2)
    # ==========================================
    slide8 = add_slide_frame(prs.slides.add_slide(slide_layout), "Key Contributions", 8)
    bullets8 = [
        "C2: Synthesis of Multi-Layered Hierarchical Protection Architectures for AC Microgrids.",
        "Designed communication-assisted coordination strategies for grid-tied and islanded transitions (Validated via hardware-in-the-loop).",
        "C3: Design of Physics-Informed Machine Learning classifiers (PINNs) for high-speed, robust fault discrimination in low-inertia environments."
    ]
    add_bullet_list(slide8, bullets8)

    # ==========================================
    # SLIDE 9: Thesis Organization
    # ==========================================
    slide9 = add_slide_frame(prs.slides.add_slide(slide_layout), "Thesis Organization", 9)
    bullets9 = [
        "Chapter 2: Physics-Informed Modelling of IBRs (PSCAD/EMT).",
        "Chapter 3: Protection requirements for Active Distribution Networks.",
        "Chapter 4: Algorithms for Multi-Functional Adaptive Interconnection Protection (Focus on C1).",
        "Chapter 5: Adaptive protection architectures for microgrids (Focus on C2).",
        "Chapter 6: AI/ML-based advanced protection utilizing PINNs (Focus on C3).",
        "Chapter 7: Conclusions and Future Work."
    ]
    add_bullet_list(slide9, bullets9)

    # ==========================================
    # SLIDE 10: Publications (Journal & Book Chapter)
    # ==========================================
    slide10 = add_slide_frame(prs.slides.add_slide(slide_layout), "List of Publications (1/2)", 10)
    bullets10 = [
        "Refereed Journals:",
        "• A. V. Eluvathingal and K. S. Swarup, 'A novel solution for overcoming conflicts between line protection and LVRT in MV Solar PV interconnections,' IEEE Access, 2024.",
        "Book Chapter:",
        "• A. V. Eluvathingal and K. S. Swarup. 'Protection scheme for smart distribution networks with inverter interfaced renewable power generating sources.' Springer Singapore, 2020."
    ]
    add_bullet_list(slide10, bullets10, top=Inches(1.5), left=Inches(0.5), width=Inches(12.333))
    # Adjust font size for this specific slide if needed
    for shape in slide10.shapes:
        if shape.has_text_frame and shape.text_frame.text.startswith("Refereed"):
            for p in shape.text_frame.paragraphs:
                p.font.size = Pt(20)

    # ==========================================
    # SLIDE 11: Publications (Conferences)
    # ==========================================
    slide11 = add_slide_frame(prs.slides.add_slide(slide_layout), "List of Publications (2/2)", 11)
    bullets11 = [
        "Selected Conference Presentations:",
        "• 'An Intelligent Methodology to Improve Distribution System Operational Parameters Utilizing Smart Inverter Functionalities of PV Sources,' RPG 2018, Denmark.",
        "• 'An interface protection relay for networked microgrids with inverter based sources,' APPEEC 2017.",
        "• 'Instantaneous symmetrical components based microgrid interface protection relay,' ICPS 2017.",
        "• 'Impact of Active Network Management Scheme in Fault Protection Design and Network Operation for Islanded Power System,' ISGT Asia 2018."
    ]
    add_bullet_list(slide11, bullets11, top=Inches(1.5), left=Inches(0.5), width=Inches(12.333))
    for shape in slide11.shapes:
        if shape.has_text_frame and shape.text_frame.text.startswith("Selected"):
            for p in shape.text_frame.paragraphs:
                p.font.size = Pt(18)

    # ==========================================
    # SLIDE 12: Conclusion
    # ==========================================
    slide12 = add_slide_frame(prs.slides.add_slide(slide_layout), "Summary & Next Steps", 12)
    bullets12 = [
        "Addressed the critical protection gaps in modern, low-inertia, IBR-dominated microgrids.",
        "Developed and validated adaptive algorithms resolving conflicts between legacy relays and LVRT mandates.",
        "Proposed a robust PINN architecture for future autonomous, self-healing grids.",
        "Thank You. Questions?"
    ]
    add_bullet_list(slide12, bullets12)

    # --- Save Presentation ---
    prs.save('TSA_Meeting_Eluvathingal.pptx')
    print("Presentation saved successfully as 'TSA_Meeting_Eluvathingal.pptx'")

if __name__ == '__main__':
    create_presentation()